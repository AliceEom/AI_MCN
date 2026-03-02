from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "slides"
ASSETS_DIR = SLIDES_DIR / "assets"
DATA_DIR = ROOT / "data"
SUMMARY_JSON = ROOT / "artifacts" / "reports" / "presentation_summary_boj.json"
OUT_PPTX = SLIDES_DIR / "AI_MCN_Final_Presentation_EN.pptx"


def add_title(slide, title: str, subtitle: str | None = None, accent: bool = True) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = RGBColor(15, 44, 89)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(11.8), Inches(0.68))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.name = "Aptos"
        sp.font.color.rgb = RGBColor(74, 95, 128)

    if accent:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.92), Inches(2.9), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(26, 125, 188)
        line.line.fill.background()


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    font_size: int = 21,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(28, 45, 70)
        p.space_after = Pt(8)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    title_size: int = 18,
    body_size: int = 15,
) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(242, 247, 252)
    rect.line.color.rgb = RGBColor(205, 220, 238)
    rect.line.width = Pt(1.25)

    title_box = slide.shapes.add_textbox(Inches(left + 0.20), Inches(top + 0.12), Inches(width - 0.30), Inches(0.38))
    ttf = title_box.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(title_size)
    tp.font.color.rgb = RGBColor(19, 58, 102)
    tp.font.name = "Aptos"

    add_bullets(
        slide,
        left=left + 0.18,
        top=top + 0.48,
        width=width - 0.28,
        height=height - 0.58,
        bullets=bullets,
        font_size=body_size,
    )


def add_two_column_bullets(
    slide,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
    top: float = 2.0,
    height: float = 4.9,
    body_size: int = 14,
) -> None:
    add_card(
        slide,
        0.68,
        top,
        5.95,
        height,
        left_title,
        left_bullets,
        title_size=17,
        body_size=body_size,
    )
    add_card(
        slide,
        6.72,
        top,
        5.95,
        height,
        right_title,
        right_bullets,
        title_size=17,
        body_size=body_size,
    )


def add_source_links(slide, links: list[tuple[str, str | None]]) -> None:
    if not links:
        return
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.03), Inches(9.3), Inches(0.34))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)

    prefix = p.add_run()
    prefix.text = "Sources: "
    prefix.font.size = Pt(8.5)
    prefix.font.color.rgb = RGBColor(110, 128, 152)
    prefix.font.name = "Aptos"

    for i, (label, url) in enumerate(links):
        if i:
            sep = p.add_run()
            sep.text = " | "
            sep.font.size = Pt(8.5)
            sep.font.color.rgb = RGBColor(130, 145, 165)
            sep.font.name = "Aptos"
        r = p.add_run()
        r.text = label
        r.font.size = Pt(8.5)
        r.font.name = "Aptos"
        if url:
            r.font.color.rgb = RGBColor(26, 125, 188)
            r.font.underline = True
            r.hyperlink.address = url
        else:
            r.font.color.rgb = RGBColor(110, 128, 152)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(10.05), Inches(7.03), Inches(2.75), Inches(0.32))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(110, 128, 152)
    p.font.name = "Aptos"


def image_path(name: str) -> Path:
    p = ASSETS_DIR / name
    if not p.exists():
        raise FileNotFoundError(f"Missing image asset: {p}")
    return p


def compute_keyword_snapshot() -> dict[str, object]:
    candidates = [
        DATA_DIR / "videos_text_ready_combined.csv",
        DATA_DIR / "videos_text_ready_demo.csv",
    ]
    source = next((p for p in candidates if p.exists()), None)
    if source is None:
        return {"source_file": "N/A", "total_videos": 0, "total_channels": 0, "keywords": {}}

    usecols = ["_channel_id", "snippet__title", "snippet__description", "snippet__tags"]
    df = pd.read_csv(source, usecols=usecols, low_memory=False)
    text = (
        df["snippet__title"].fillna("")
        + " "
        + df["snippet__description"].fillna("")
        + " "
        + df["snippet__tags"].fillna("")
    ).str.lower()

    keywords = ["sunscreen", "spf", "k-beauty", "beauty of joseon", "cerave"]
    out: dict[str, dict[str, int]] = {}
    for kw in keywords:
        m = text.str.contains(kw, regex=False)
        out[kw] = {
            "videos": int(m.sum()),
            "channels": int(df.loc[m, "_channel_id"].astype(str).nunique()),
        }

    return {
        "source_file": source.name,
        "total_videos": int(len(df)),
        "total_channels": int(df["_channel_id"].astype(str).nunique()),
        "keywords": out,
    }


def _fmt_kw(kw_snapshot: dict[str, object], keyword: str) -> str:
    row = (kw_snapshot.get("keywords") or {}).get(keyword, {})
    return f"{keyword}: {int(row.get('videos', 0)):,} videos | {int(row.get('channels', 0)):,} channels"


def build() -> Path:
    if not SUMMARY_JSON.exists():
        raise FileNotFoundError(f"Missing summary file: {SUMMARY_JSON}")
    summary = json.loads(SUMMARY_JSON.read_text(encoding="utf-8"))

    dataset = summary.get("dataset", {})
    roi = summary.get("roi", {})
    ml_table = summary.get("ml_table", [])
    top10 = summary.get("top10", [])
    kw_snapshot = compute_keyword_snapshot()

    best_model = summary.get("ml_best_model", "N/A")
    baseline_rmse = next((x.get("rmse_mean", 0.0) for x in ml_table if x.get("model") == "BaselineMedian"), 0.0)
    best_rmse = next((x.get("rmse_mean", 0.0) for x in ml_table if x.get("model") == best_model), 0.0)
    gain_pct = ((baseline_rmse - best_rmse) / baseline_rmse * 100.0) if baseline_rmse else 0.0

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title and team
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "AI-MCN: Replacing Traditional MCN Matching with an AI Decision System",
        "MSIS 521 Course Project | 15-minute Presentation",
    )
    add_bullets(
        slide,
        0.72,
        2.2,
        11.9,
        2.5,
        [
            "Client case: Beauty of Joseon (BOJ), U.S. sunscreen campaign",
            "Goal: help brands find best-fit influencers directly with transparent AI logic",
            "Team: [Add names]",
        ],
        font_size=23,
    )
    add_source_links(
        slide,
        [
            ("GitHub repo", "https://github.com/AliceEom/AI_MCN"),
            ("EMARKETER 2025", "https://www.emarketer.com/press-releases/us-influencer-marketing-spending-will-surpass-10-billion-in-2025/"),
        ],
    )
    add_footer(slide, "Section 1 | Title and Team")

    # 2. Agenda and timing
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Agenda and Timing (Professor Guide)")
    add_bullets(
        slide,
        0.82,
        2.1,
        12.0,
        4.8,
        [
            "1) Title and team (1 min)",
            "2) Business context and problem (2-3 min)",
            "3) Data and features (1-3 min)",
            "4) AI/ML approach (3-4 min)",
            "5) Prototype demo (3-4 min)",
            "6) Impact, limitations, next steps (2-3 min)",
            "7) How we used AI tools (1-2 min)",
        ],
        font_size=21,
    )
    add_source_links(
        slide,
        [
            ("MSIS 521 assignment/rubric (internal PDFs)", None),
        ],
    )
    add_footer(slide, "Section 2 | Roadmap")

    # 3. Business context
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Business Context and Problem")
    add_card(
        slide,
        0.65,
        2.0,
        5.95,
        4.8,
        "Current MCN-Style Workflow",
        [
            "MCNs aggregate influencer data and broker brand-creator matching.",
            "Brands pay agency/management fees plus campaign budget.",
            "Matching logic is often hard to audit for brand teams.",
        ],
    )
    add_card(
        slide,
        6.72,
        2.0,
        5.95,
        4.8,
        "Decision We Improve",
        [
            "From manual or popularity-led creator selection",
            "To evidence-based, product-specific creator ranking",
            "Goal: reduce expensive mismatch risk before launch",
        ],
    )
    add_source_links(
        slide,
        [
            ("YouTube Content Manager", "https://support.google.com/youtube/answer/106934?hl=en-EN"),
            ("YouTube MCN setup", "https://support.google.com/youtube/answer/7296308?hl=en-GB"),
            ("IAB Creator Economy", "https://www.iab.com/news/creator-economy-ad-spend-to-reach-37-billion-in-2025-growing-4x-faster-than-total-media-industry-according-to-iab/"),
        ],
    )
    add_footer(slide, "Section 3 | Context and Decision")

    # 4. Why now / market evidence
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why Now: Market Evidence")
    add_bullets(
        slide,
        0.8,
        2.0,
        12.0,
        4.7,
        [
            "US influencer marketing spend is projected at $10.52B in 2025 (EMARKETER, Mar 13, 2025).",
            "YouTube influencer use among US marketers is projected to exceed 50% in 2025 (EMARKETER).",
            "Global beauty grew 7.3% YoY in 2025; US beauty sales are 41% e-commerce (NIQ, Feb 25, 2025).",
            "McKinsey projects core beauty segments to reach about $590B by 2030, with skincare as the largest segment (Aug 28, 2025).",
        ],
        font_size=18,
    )
    add_source_links(
        slide,
        [
            ("EMARKETER", "https://www.emarketer.com/press-releases/us-influencer-marketing-spending-will-surpass-10-billion-in-2025/"),
            ("NIQ Beauty 2025", "https://nielseniq.com/global/en/news-center/2025/niq-reports-7-3-year-over-year-value-growth-in-global-beauty-sector/"),
            ("McKinsey 2025", "https://www.mckinsey.com/industries/consumer-packaged-goods/our-insights/a-close-look-at-the-global-beauty-industry-in-2025"),
        ],
    )
    add_footer(slide, "Section 3 | External Industry Context")

    # 5. Why beauty and BOJ
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why Beauty and Why BOJ")
    add_card(
        slide,
        0.65,
        2.0,
        5.95,
        4.75,
        "Case Fit",
        [
            "Beauty creator marketing is large and highly content-driven.",
            "BOJ has visible social momentum and a product-specific campaign story.",
            "BOJ U.S. expansion + sunscreen focus make a practical demo scenario.",
            "(Fashionista coverage, Nov 13, 2024)",
        ],
    )
    add_card(
        slide,
        6.72,
        2.0,
        5.95,
        4.75,
        "Dataset Feasibility Signals",
        [
            f"Source file: {kw_snapshot.get('source_file', 'N/A')} (n={int(kw_snapshot.get('total_videos', 0)):,} videos)",
            _fmt_kw(kw_snapshot, "sunscreen"),
            _fmt_kw(kw_snapshot, "spf"),
            _fmt_kw(kw_snapshot, "beauty of joseon"),
            _fmt_kw(kw_snapshot, "cerave"),
        ],
    )
    add_source_links(
        slide,
        [
            ("Fashionista BOJ", "https://fashionista.com/2024/11/beauty-of-joseon-k-beauty-skin-care-us-launch-strategy"),
            ("BOJ Rice Wonderland", "https://beautyofjoseon.com/pages/rice-wonderland"),
            ("FDA SIA", "https://www.fda.gov/drugs/guidance-compliance-regulatory-information/sunscreen-innovation-act-sia"),
        ],
    )
    add_footer(slide, "Section 3 | Client and Case Selection")

    # 6. Data and features
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Data and Features")
    add_card(
        slide,
        0.7,
        2.0,
        6.0,
        4.7,
        "Data Used",
        [
            "Team-collected YouTube API exports (videos/comments/channel fields)",
            f"Videos analyzed (full run): {int(dataset.get('videos_analyzed', 0)):,}",
            f"Channels scored: {int(dataset.get('channels_scored', 0)):,}",
            "Main files: videos_text_ready, comments_raw, master_prd_slim",
        ],
    )
    add_card(
        slide,
        6.75,
        2.0,
        5.9,
        4.7,
        "Feature Groups",
        [
            "Network features: centrality + community",
            "Text features: TF-IDF + semantic/tone alignment",
            "Behavioral features: views/likes/comments/engagement",
            "Reliability features: evidence score + credibility multiplier",
        ],
    )
    add_source_links(
        slide,
        [
            ("YouTube Data API", "https://developers.google.com/youtube/v3"),
            ("Project data files (internal)", None),
        ],
    )
    add_footer(slide, "Section 4 | Data and Feature Engineering")

    # 7. Preprocessing and EDA
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Preprocessing and EDA")
    slide.shapes.add_picture(str(image_path("community_distribution_clean.png")), Inches(0.78), Inches(2.0), height=Inches(4.8))
    add_card(
        slide,
        8.0,
        2.0,
        4.5,
        4.8,
        "What we did",
        [
            "Deduplication + type normalization",
            "Beauty include filter + non-beauty noise exclusion",
            "Channel-level aggregation and recency features",
            "EDA confirmed cluster concentration -> diversity guardrail needed",
        ],
        body_size=14,
    )
    add_source_links(
        slide,
        [
            ("Preprocessing pipeline (internal)", None),
            ("EDA plots (internal)", None),
        ],
    )
    add_footer(slide, "Section 4 | Data Understanding")

    # 8. AI/ML approach
    slide = prs.slides.add_slide(blank)
    add_title(slide, "AI/ML Approach (Class-Learned Methods Focus)")
    add_bullets(
        slide,
        0.75,
        2.0,
        12.0,
        4.9,
        [
            "1) Social Network Analysis: creator graph, centrality, communities",
            "2) Text Analytics: TF-IDF relevance + semantic/tone enrichment",
            "3) Supervised Regression: Linear, LASSO, Ridge, CART, RandomForest, LightGBM",
            "4) Validation: GroupKFold(5) by channel to reduce leakage",
            "5) Explainability: SHAP for feature contribution",
            "Extensions shown in demo: ROI simulator, strategy and memo generation",
        ],
        font_size=18,
    )
    add_source_links(
        slide,
        [
            ("GroupKFold", "https://scikit-learn.org/stable/modules/generated/sklearn.model_selection.GroupKFold.html"),
            ("LightGBM", "https://lightgbm.readthedocs.io/en/latest/"),
            ("SHAP", "https://shap.readthedocs.io/en/latest/"),
        ],
    )
    add_footer(slide, "Section 5 | Methods")

    # 9. Why this method mix
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why This Method Mix")
    add_card(
        slide,
        0.65,
        2.0,
        5.95,
        4.8,
        "Why not follower-count only",
        [
            "Misses campaign-language fit",
            "Misses network position and community overlap",
            "Misses quality/reliability signals",
            "High chance of expensive creator mismatch",
        ],
    )
    add_card(
        slide,
        6.72,
        2.0,
        5.95,
        4.8,
        "Hybrid ranking logic",
        [
            "Network influence + text fit + engagement potential",
            "Reliability multiplier to penalize low-signal channels",
            "Diversity constraint for healthier final shortlist",
            "More robust than single-signal ranking",
        ],
    )
    add_source_links(
        slide,
        [
            ("FTC Influencers", "https://www.ftc.gov/influencers"),
            ("FTC Endorsement Guides", "https://www.ftc.gov/business-guidance/resources/ftcs-endorsement-guides"),
            ("YouTube MCN reference", "https://support.google.com/youtube/answer/7296308?hl=en-GB"),
        ],
    )
    add_footer(slide, "Section 5 | Method Choice Rationale")

    # 10. Model results
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Model Results and Validation")
    slide.shapes.add_picture(str(image_path("model_benchmark_rmse.png")), Inches(0.78), Inches(2.0), height=Inches(4.85))
    add_card(
        slide,
        8.0,
        2.0,
        4.45,
        4.85,
        f"Best model: {best_model}",
        [
            f"Best RMSE: {best_rmse:.5f}",
            f"Baseline RMSE: {baseline_rmse:.5f}",
            f"Relative RMSE reduction: {gain_pct:.1f}%",
            "Interpretation: meaningful uplift over naive baseline",
        ],
    )
    add_source_links(
        slide,
        [
            ("ML CV outputs (internal)", None),
            ("LightGBM docs", "https://lightgbm.readthedocs.io/en/latest/"),
        ],
    )
    add_footer(slide, "Section 5 | Validation Output")

    # 11. Demo and use cases
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Prototype Demo: Input to Output")
    slide.shapes.add_picture(str(image_path("top10_final_scores.png")), Inches(7.0), Inches(2.1), width=Inches(5.5))
    add_bullets(
        slide,
        0.8,
        2.0,
        6.0,
        4.9,
        [
            "Live flow: campaign input -> Top-N shortlist -> rationale/risk",
            "Interactive tabs: Network, Text Intelligence, ML, ROI",
            "Use case A: BOJ sunscreen launch planning",
            "Use case B: CeraVe benchmark for shortlist calibration",
            "Decision output: ranked creators + risk notes + memo export",
        ],
        font_size=17,
    )
    add_source_links(
        slide,
        [
            ("Streamlit docs", "https://docs.streamlit.io/"),
            ("Prototype app/pipeline (internal)", None),
        ],
    )
    add_footer(slide, "Section 6 | Live Demo")

    # 12. Impact
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Business Impact")
    add_card(
        slide,
        0.68,
        2.0,
        6.0,
        4.8,
        "How this helps the client",
        [
            "Faster influencer shortlisting and campaign planning",
            "Transparent recommendation logic for stakeholder trust",
            "Reduced mismatch risk before spending budget",
            "Scenario-based ROI planning and expectation setting",
        ],
    )
    add_card(
        slide,
        6.75,
        2.0,
        5.9,
        4.8,
        "Decisions improved",
        [
            "Who to prioritize in outreach",
            "How many creators to activate by objective",
            "How much budget risk is acceptable",
            f"Base ROI scenario in run: {float(roi.get('roas', 0)):.2f}x expected ROAS",
        ],
    )
    add_source_links(
        slide,
        [
            ("IAB Creator Economy", "https://www.iab.com/news/creator-economy-ad-spend-to-reach-37-billion-in-2025-growing-4x-faster-than-total-media-industry-according-to-iab/"),
            ("YouTube U.S. Impact 2024", "https://blog.youtube/news-and-events/2024-us-youtube-impact-report/"),
        ],
    )
    add_footer(slide, "Section 7 | Business Value")

    # 13. Limitations and next steps
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Limitations, Ethics, and Next Steps")
    add_card(
        slide,
        0.7,
        2.0,
        5.95,
        4.8,
        "Current limitations",
        [
            "Pre-collected YouTube dataset (not live cross-platform feed)",
            "ROI is scenario estimate, not causal proof",
            "Remaining cluster concentration in skincare niches",
            "Need larger-scale operational monitoring",
        ],
    )
    add_card(
        slide,
        6.75,
        2.0,
        5.95,
        4.8,
        "Ethics and roadmap",
        [
            "Bias checks: degree-overlap diagnostics + evidence penalty",
            "Privacy: no sensitive personal profiling in prototype",
            "Next: live connectors + fairness constraints + pilot calibration",
            "Target: weekly decision-support workflow for brand teams",
        ],
    )
    add_source_links(
        slide,
        [
            ("FTC Influencer Disclosures", "https://www.ftc.gov/influencers"),
            ("SAFE Sunscreen (House)", "https://www.congress.gov/bill/119th-congress/house-bill/3686"),
            ("FDA SIA", "https://www.fda.gov/drugs/guidance-compliance-regulatory-information/sunscreen-innovation-act-sia"),
        ],
    )
    add_footer(slide, "Section 7 | Caveats and Roadmap")

    # 14. AI tools usage
    slide = prs.slides.add_slide(blank)
    add_title(slide, "How We Used AI Tools")
    add_card(
        slide,
        0.7,
        2.0,
        5.95,
        4.75,
        "Where AI tools helped",
        [
            "Ideation and scope refinement",
            "Implementation acceleration and debugging",
            "Documentation and memo drafting",
            "Slide polishing and presentation prep",
        ],
    )
    add_card(
        slide,
        6.75,
        2.0,
        5.95,
        4.75,
        "Human responsibility",
        [
            "Problem framing and business assumptions",
            "Model choice, validation, and guardrail design",
            "Result interpretation and recommendation decisions",
            "Final accountability stayed with the team",
        ],
    )
    add_source_links(
        slide,
        [
            ("Project commits/artifacts (internal)", None),
            ("OpenAI platform docs", "https://platform.openai.com/docs/overview"),
        ],
    )
    add_footer(slide, "Section 8 | AI Tooling Disclosure")

    # 15. References and Q&A
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Conclusion, References, and Q&A")
    add_bullets(
        slide,
        0.82,
        2.0,
        12.0,
        3.9,
        [
            "Conclusion: AI-MCN can reduce manual MCN dependence and improve influencer decision quality.",
            "EMARKETER (Mar 13, 2025): US influencer spend and YouTube marketer adoption",
            "NIQ (Feb 25, 2025): beauty growth and e-commerce/social commerce indicators",
            "McKinsey (Aug 28, 2025): beauty market outlook to 2030",
            "Fashionista (Nov 13, 2024): BOJ virality and U.S. market context",
            "Technical refs: YouTube API, scikit-learn, LightGBM, SHAP, Streamlit",
        ],
        font_size=16,
    )
    q = slide.shapes.add_textbox(Inches(0.82), Inches(5.85), Inches(11.8), Inches(0.65))
    qp = q.text_frame.paragraphs[0]
    qp.text = "Thank you. Q&A"
    qp.font.size = Pt(30)
    qp.font.bold = True
    qp.font.name = "Aptos Display"
    qp.font.color.rgb = RGBColor(16, 86, 150)
    qp.alignment = PP_ALIGN.CENTER
    add_source_links(
        slide,
        [
            ("EMARKETER", "https://www.emarketer.com/press-releases/us-influencer-marketing-spending-will-surpass-10-billion-in-2025/"),
            ("NIQ", "https://nielseniq.com/global/en/news-center/2025/niq-reports-7-3-year-over-year-value-growth-in-global-beauty-sector/"),
            ("McKinsey", "https://www.mckinsey.com/industries/consumer-packaged-goods/our-insights/a-close-look-at-the-global-beauty-industry-in-2025"),
        ],
    )
    add_footer(slide, "AI-MCN | MSIS 521")

    # 16. Backup: model performance details
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Model Performance Snapshot")
    slide.shapes.add_picture(str(image_path("model_benchmark_rmse.png")), Inches(0.78), Inches(2.0), height=Inches(4.8))
    add_card(
        slide,
        7.95,
        2.0,
        4.5,
        4.8,
        "BOJ Full-Run Metrics",
        [
            f"Best model: {best_model}",
            f"Best RMSE: {best_rmse:.5f}",
            f"Baseline RMSE: {baseline_rmse:.5f}",
            f"Relative RMSE reduction: {gain_pct:.1f}%",
            "Best model selected by 5-fold GroupKFold RMSE.",
            "BaselineMedian retained as reference in output table.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("ml_cv_results.csv (internal)", None),
            ("presentation_summary_boj.json (internal)", None),
            ("GroupKFold docs", "https://scikit-learn.org/stable/modules/generated/sklearn.model_selection.GroupKFold.html"),
        ],
    )
    add_footer(slide, "Backup 16 | Model Results Detail")

    # 17. Backup: top-match anatomy
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Top-Match Recommendation Anatomy")
    add_two_column_bullets(
        slide,
        "What Each Recommendation Card Includes",
        [
            "Final Match Score (current ranking score in the tab)",
            "Signal breakdown: SNA, TF-IDF, Semantic, Tone, Engagement, ML",
            "Score controls: Base score, reliability multiplier, community id",
            "Evidence label (High/Medium/Low) and fit label (Very Strong to Exploratory)",
            "Channel snapshot, keyword summary, best/recent videos, recent comments",
            "Red flags and rationale text for reviewer transparency",
        ],
        "How to Read It for Decision-Making",
        [
            "Start with Final Match Score and Evidence together.",
            "Check reliability multiplier before trusting high base scores.",
            "Review red flags for exclusions, low semantic fit, low activity.",
            "Use video/comment context to verify qualitative brand fit.",
            "Compare component mix by objective (awareness vs conversion).",
            "Use diversity settings to avoid over-concentrated shortlist.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("Top Matches tab code (app.py)", None),
            ("channel_details.py (internal)", None),
            ("ranking.py (internal)", None),
        ],
    )
    add_footer(slide, "Backup 17 | Recommendation Interpretation")

    # 18. Backup: network studio interpretation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Network Studio Interpretation Guide")
    add_two_column_bullets(
        slide,
        "What the Network View Represents",
        [
            "Node = creator channel; edge = shared-tag relationship.",
            "Edge weight = count of shared tags above threshold.",
            "Node color = community cluster id.",
            "Node size increases with final score in graph view.",
            "Over-common tags are dropped to reduce graph noise.",
            "Micro communities can be separated as community -1.",
        ],
        "How Bias Diagnostics Are Used",
        [
            "Degree Top Overlap compares centrality-only vs hybrid ranking.",
            "Lower overlap suggests reduced popularity-only bias.",
            "Unique communities in Top-N reflects portfolio diversity.",
            "Community distribution chart reveals concentration risk.",
            "Min edge / node sliders help stress-test shortlist stability.",
            "Use diagnostics before final outreach decisions.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("network_scoring.py (internal)", None),
            ("orchestrator.py bias report (internal)", None),
            ("Network Studio UI (app.py)", None),
        ],
    )
    add_footer(slide, "Backup 18 | Network Diagnostics")

    # 19. Backup: text intelligence interpretation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Text Intelligence Interpretation Guide")
    add_two_column_bullets(
        slide,
        "Charts and Metrics",
        [
            "Text Match Map: X=TF-IDF, Y=Semantic, color=Evidence, size~views.",
            "Top Terms: frequent words in top candidate text corpus.",
            "Keyword Coverage: share of channels containing each tracked keyword.",
            "Leaderboard: TF-IDF, Semantic, Tone, Final score side-by-side.",
            "Campaign words from audience/USP/must-keywords shape query vector.",
            "Exclude words can directly lower relevance scores.",
        ],
        "How to Use in Client Reviews",
        [
            "Upper-right map region = strongest language + meaning fit.",
            "High TF-IDF but low semantic means superficial keyword overlap.",
            "Low must-keyword coverage indicates query or data-gap issue.",
            "Term chart validates whether shortlist matches campaign language.",
            "Use with Top-Match evidence labels for safer final picks.",
            "Document rationale from this tab in client memo.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("text_scoring.py (internal)", None),
            ("semantic_enrichment.py (internal)", None),
            ("Text Intelligence UI (app.py)", None),
        ],
    )
    add_footer(slide, "Backup 19 | Text Diagnostics")

    # 20. Backup: ROI lab interpretation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: ROI Lab Interpretation Guide")
    add_two_column_bullets(
        slide,
        "ROI Simulator Mechanics",
        [
            "Inputs: Budget, CPM, CTR, CVR, AOV",
            "Impressions = (Budget / CPM) * 1000",
            "Clicks = Impressions * CTR",
            "Conversions = Clicks * CVR",
            "Revenue = Conversions * AOV",
            "ROAS = Revenue / Budget (plus low/high range)",
        ],
        "Business Use and Caveats",
        [
            "Use for scenario planning, not causal proof.",
            "Budget-sensitivity chart shows efficiency vs scale tradeoff.",
            "Useful for pre-campaign expectation setting with finance teams.",
            "Can compare assumptions across shortlist compositions.",
            "Actual performance still depends on creative and execution quality.",
            "Should be recalibrated with real campaign outcomes.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("roi_simulation.py (internal)", None),
            ("ROI tab UI (app.py)", None),
            ("IAB creator spend context", "https://www.iab.com/news/creator-economy-ad-spend-to-reach-37-billion-in-2025-growing-4x-faster-than-total-media-industry-according-to-iab/"),
        ],
    )
    add_footer(slide, "Backup 20 | ROI Interpretation")

    # 21. Backup: business impact depth
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Business Impact Detail")
    add_two_column_bullets(
        slide,
        "Expected Impact for Brand Teams",
        [
            "Shortens creator research-to-shortlist cycle.",
            "Converts opaque creator selection into transparent scorecards.",
            "Adds reliability penalties to reduce false-positive picks.",
            "Supports benchmark-based discussions with stakeholders.",
            "Enables repeatable weekly operating rhythm.",
            "Improves auditability for campaign post-mortems.",
        ],
        "Why the Timing Is Urgent",
        [
            "Creator budget scale is increasing rapidly.",
            "Teams still report time and ROI measurement friction.",
            "Beauty category adds extra pressure via high social influence.",
            "Higher CAC increases cost of poor creator matching.",
            "Brands need decision systems, not one-off creator lists.",
            "AI-MCN framing reduces dependency on opaque intermediaries.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("IAB Creator Economy 2025", "https://www.iab.com/news/creator-economy-ad-spend-to-reach-37-billion-in-2025-growing-4x-faster-than-total-media-industry-according-to-iab/"),
            ("EMARKETER 2025", "https://www.emarketer.com/press-releases/us-influencer-marketing-spending-will-surpass-10-billion-in-2025/"),
            ("Sprout workload findings", "https://sproutsocial.com/insights/data/social-media-teams-wear-many-hats-to-reach-business-goals/"),
        ],
    )
    add_footer(slide, "Backup 21 | Impact Detail")

    # 22. Backup: limitations and governance
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Limitations and Governance Detail")
    add_two_column_bullets(
        slide,
        "Current Technical Constraints",
        [
            "Dataset is pre-collected, not full live multi-platform ingestion.",
            "Coverage depends on available YouTube data and keyword filters.",
            "ROI output is assumption-based simulation.",
            "Semantic enrichment quality depends on text richness.",
            "Community structure may still be concentrated in sub-niches.",
            "Cloud demo mode can use reduced snapshot data.",
        ],
        "Governance and Compliance Controls",
        [
            "FTC disclosure and endorsement compliance must be enforced.",
            "Claim safety checks needed for product-specific messaging.",
            "Brand safety workflow should include human reviewer checkpoints.",
            "Bias checks should be monitored as data evolves.",
            "Recommendation logs should be retained for auditability.",
            "Legal/marketing sign-off needed before launch decisions.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("FTC Influencers", "https://www.ftc.gov/influencers"),
            ("FTC Endorsement Guides", "https://www.ftc.gov/business-guidance/resources/ftcs-endorsement-guides"),
            ("FDA SIA", "https://www.fda.gov/drugs/guidance-compliance-regulatory-information/sunscreen-innovation-act-sia"),
        ],
    )
    add_footer(slide, "Backup 22 | Limitations & Compliance")

    # 23. Backup: roadmap
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: 90-Day Operational Roadmap")
    add_two_column_bullets(
        slide,
        "Phase 1 (0-30 Days)",
        [
            "Lock weekly pipeline schedule and dashboard review cadence.",
            "Define reviewer checklist for low-evidence channels.",
            "Standardize export package for brand decision meetings.",
            "Set benchmark brands by category for context comparison.",
            "Track turnaround time and reviewer agreement metrics.",
        ],
        "Phase 2-3 (31-90 Days)",
        [
            "Add live connectors and larger benchmark panel coverage.",
            "Run pilot with real outreach and campaign outcomes.",
            "Compare predicted vs realized metrics and recalibrate model.",
            "Add fairness constraints and stronger monitoring alerts.",
            "Finalize operating playbook for in-house AI-MCN workflow.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("orchestrator.py (internal)", None),
            ("app.py workflow (internal)", None),
            ("AI_MCN_Analysis_Explainer_EN.md (internal)", None),
        ],
    )
    add_footer(slide, "Backup 23 | Roadmap")

    # 24. Backup: AI tool usage detail
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: AI Tool Usage Detail")
    add_two_column_bullets(
        slide,
        "Where AI Tools Accelerated Work",
        [
            "Research synthesis and structuring references for slides.",
            "Implementation support across data, ML, and UI modules.",
            "Debugging support for pipeline/runtime issues.",
            "Documentation drafting and bilingual explainer production.",
            "Presentation refinement and narrative iteration.",
        ],
        "Human-Owned Responsibilities",
        [
            "Client problem framing and objective definition.",
            "Method selection and validation logic design.",
            "Quality control on model outputs and interpretation.",
            "Risk/compliance framing and stakeholder communication.",
            "Final recommendation accountability remained with team.",
        ],
        body_size=13,
    )
    add_source_links(
        slide,
        [
            ("OpenAI docs", "https://platform.openai.com/docs/overview"),
            ("Internal repo artifacts", None),
        ],
    )
    add_footer(slide, "Backup 24 | AI Tooling Detail")

    # 25. Backup: reference library
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Backup: Research Library for Team Follow-up")
    add_bullets(
        slide,
        0.78,
        2.0,
        12.0,
        4.9,
        [
            "Market spend: EMARKETER (Mar 13, 2025), IAB (Nov 20, 2025), Goldman Sachs creator economy outlook.",
            "Operational pain: Sprout Social workload findings; Linqia 2023/2026 benchmark reports.",
            "Beauty context: NIQ global beauty growth + K-beauty landscape; McKinsey 2025 beauty outlook.",
            "Client context: BOJ official campaign pages + public U.S. strategy coverage.",
            "Governance: FTC influencer disclosures/endorsement guides, YouTube MCN documentation, FDA sunscreen policy context.",
            "All links are included in slide markdown and sources footer for teammate reading.",
        ],
        font_size=15,
    )
    add_source_links(
        slide,
        [
            ("EMARKETER", "https://www.emarketer.com/press-releases/us-influencer-marketing-spending-will-surpass-10-billion-in-2025/"),
            ("IAB", "https://www.iab.com/news/creator-economy-ad-spend-to-reach-37-billion-in-2025-growing-4x-faster-than-total-media-industry-according-to-iab/"),
            ("NIQ", "https://nielseniq.com/global/en/news-center/2025/niq-reports-7-3-year-over-year-value-growth-in-global-beauty-sector/"),
            ("McKinsey", "https://www.mckinsey.com/industries/consumer-packaged-goods/our-insights/a-close-look-at-the-global-beauty-industry-in-2025"),
        ],
    )
    add_footer(slide, "Backup 25 | References")

    prs.save(str(OUT_PPTX))
    return OUT_PPTX


if __name__ == "__main__":
    out = build()
    print(out)
