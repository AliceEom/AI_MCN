from __future__ import annotations

import os
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.usecase_planner import build_usecase_plan

ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "slides"
ASSETS_DIR = SLIDES_DIR / "assets"
REPORTS_DIR = ROOT / "artifacts" / "reports"
OUT_PPTX = SLIDES_DIR / "AI_MCN_Demo_DeepDive_EN.pptx"
DEMO_URL = os.getenv("AI_MCN_DEMO_URL", "https://YOUR-STREAMLIT-APP-URL")


# ---------- Basic drawing helpers ----------
def _set_slide_bg_white(slide) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    _set_slide_bg_white(slide)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = RGBColor(17, 56, 101)

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(1.06), Inches(12.2), Inches(0.55))
        stf = sb.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.name = "Aptos"
        sp.font.color.rgb = RGBColor(80, 96, 118)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.58), Inches(3.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(35, 116, 203)
    line.line.fill.background()


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(10.1), Inches(7.0), Inches(2.8), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.name = "Aptos"
    p.font.color.rgb = RGBColor(112, 126, 146)


def add_source(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(9.7), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Source: {text}"
    p.font.size = Pt(8.5)
    p.font.name = "Aptos"
    p.font.color.rgb = RGBColor(110, 128, 150)


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    font_size: int = 17,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(28, 45, 70)
        p.space_after = Pt(5)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(244, 248, 253)
    card.line.color.rgb = RGBColor(210, 223, 240)
    card.line.width = Pt(1.15)

    tb = slide.shapes.add_textbox(Inches(left + 0.16), Inches(top + 0.12), Inches(width - 0.28), Inches(0.45))
    ttf = tb.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(title_size)
    tp.font.name = "Aptos"
    tp.font.color.rgb = RGBColor(19, 58, 102)

    add_bullets(slide, left + 0.16, top + 0.50, width - 0.28, height - 0.60, bullets, font_size=body_size)


def add_link_box(slide, left: float, top: float, width: float, height: float, label: str, url: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(225, 241, 255)
    shape.line.color.rgb = RGBColor(35, 116, 203)
    shape.line.width = Pt(1.6)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Aptos"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(17, 56, 101)
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].hyperlink.address = url


def add_optional_image(slide, name: str, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    path = ASSETS_DIR / name
    if not path.exists():
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


# ---------- Data preparation ----------
def _load_inputs() -> tuple[pd.DataFrame, pd.DataFrame, dict, str]:
    top = pd.read_csv(REPORTS_DIR / "top10_beauty_of_joseon.csv")
    scored = pd.read_csv(REPORTS_DIR / "scored_channels_beauty_of_joseon.csv")
    summary = pd.read_json(REPORTS_DIR / "presentation_summary_boj.json", typ="series").to_dict()
    memo = (REPORTS_DIR / "memo_beauty_of_joseon.md").read_text(encoding="utf-8")
    return top, scored, summary, memo


def _build_channel_plan(top: pd.DataFrame, roi: dict) -> pd.DataFrame:
    return build_usecase_plan(top, roi).plan_df


def _risk_lines(memo_text: str, max_items: int = 6) -> list[str]:
    lines = memo_text.splitlines()
    in_risk = False
    out: list[str] = []
    for ln in lines:
        t = ln.strip()
        if t.startswith("## Risk Flags"):
            in_risk = True
            continue
        if in_risk and t.startswith("## "):
            break
        if in_risk and t.startswith("-"):
            out.append(t.lstrip("-").strip())
            if len(out) >= max_items:
                break
    return out if out else ["No explicit risk flags found in memo."]


def _line_for_channel(r: pd.Series) -> str:
    return (
        f"{r['channel_title']} | {r['activation_tier']} | {r['content_concept']} | "
        f"Budget~${int(r['plan_budget_usd']):,} | Conv~{int(r['plan_conversions']):,}"
    )


def _comma_join(values: list[str]) -> str:
    vals = [str(v).strip() for v in values if str(v).strip()]
    if not vals:
        return "None in current Top-10"
    return ", ".join(vals)


def _role_split_lines(plan: pd.DataFrame) -> list[str]:
    if plan.empty:
        return ["Role split unavailable (empty plan table)."]

    awareness = plan[plan["activation_tier"] == "Awareness-focused"]["channel_title"].astype(str).tolist()
    conversion = plan[plan["activation_tier"] == "Conversion-focused"]["channel_title"].astype(str).tolist()
    support = plan[~plan["activation_tier"].isin(["Awareness-focused", "Conversion-focused"])]["channel_title"].astype(str).tolist()
    return [
        f"Scale/Awareness: {_comma_join(awareness)}",
        f"Conversion/Fit: {_comma_join(conversion)}",
        f"Education/Support: {_comma_join(support)}",
        "Use different creative objective by role, not one message for all creators.",
    ]


def _concept_mapping_lines(plan: pd.DataFrame) -> list[str]:
    if plan.empty:
        return ["Concept mapping unavailable (empty plan table)."]
    c1 = plan[plan["content_concept"].astype(str).str.contains("Concept 1", na=False)]["channel_title"].astype(str).tolist()
    c2 = plan[plan["content_concept"].astype(str).str.contains("Concept 2", na=False)]["channel_title"].astype(str).tolist()
    c3 = plan[plan["content_concept"].astype(str).str.contains("Concept 3", na=False)]["channel_title"].astype(str).tolist()
    return [
        f"Concept 1 (Daily Routine): {_comma_join(c1)}",
        f"Concept 2 (Results Comparison): {_comma_join(c2)}",
        f"Concept 3 (Q&A Hook): {_comma_join(c3)}",
        "Assign creator objective before budget allocation.",
    ]


# ---------- Deck build ----------
def build() -> Path:
    top, scored, summary, memo = _load_inputs()

    roi = summary.get("roi", {}) or {}
    benchmark = summary.get("benchmark", {}) or {}
    videos = int(summary.get("dataset", {}).get("videos_analyzed", len(scored)))
    channels = int(summary.get("dataset", {}).get("channels_scored", scored.get("_channel_id", pd.Series(dtype=str)).nunique()))
    best_model = str(summary.get("ml_best_model", "N/A"))

    plan = _build_channel_plan(top, roi)
    risk = _risk_lines(memo)

    top_mean = float(top["final_score"].mean()) if not top.empty else 0.0
    ev_mean = float(top["evidence_score"].mean()) if not top.empty else 0.0
    comm_n = int(top["community_id"].nunique()) if not top.empty else 0
    bench_mean = float(benchmark.get("topn_mean_score", 0.0))
    bench_gap = top_mean - bench_mean

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) link
    s = prs.slides.add_slide(blank)
    add_title(s, "Live Demo", "Click to launch the Streamlit app")
    add_link_box(s, 1.0, 2.25, 11.3, 1.15, "Open AI-MCN Demo", DEMO_URL)
    add_card(
        s,
        0.85,
        3.75,
        12.0,
        2.7,
        "What I will do in live demo",
        [
            "Explain every tab from input to export.",
            "Define each important metric in business language.",
            "Then return to slides for BOJ action plan and beauty-industry insights.",
        ],
        body_size=14,
    )
    add_source(s, "AI-MCN Streamlit prototype")
    add_footer(s, "1 / 15")

    # 2) roadmap
    s = prs.slides.add_slide(blank)
    add_title(s, "Live Demo Roadmap", "All tabs will be explained")
    add_bullets(
        s,
        0.9,
        1.95,
        6.2,
        4.9,
        [
            "1) Campaign Input",
            "2) Overview",
            "3) Top Matches",
            "4) Network Studio",
            "5) Text Intelligence",
            "6) ML Studio",
            "7) ROI Lab",
            "8) Content Strategy",
            "9) Executive Memo / Glossary / Export",
        ],
        font_size=18,
    )
    add_card(
        s,
        7.0,
        1.95,
        5.45,
        4.9,
        "Audience should watch for",
        [
            "How score is built and controlled.",
            "How reliability changes shortlist quality.",
            "How content + ROI decisions are connected.",
            "How BOJ can act immediately after this run.",
        ],
        body_size=15,
    )
    add_source(s, "AI_MCN_Demo_DeepDive_Script_EN.md")
    add_footer(s, "2 / 15")

    # 3) checkpoint slide
    s = prs.slides.add_slide(blank)
    add_title(s, "Live Demo Checkpoint", "What each tab must answer")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Decision questions",
        [
            "Who should BOJ contact first?",
            "Why are those creators recommended?",
            "How much evidence supports each recommendation?",
            "What budget plan is realistic?",
            "What content format should each creator publish?",
        ],
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Metric meanings to explain live",
        [
            "Final Score = suitability after reliability adjustment.",
            "Evidence Score = scale/activity/interaction support.",
            "Community ID = network cluster label.",
            "TF-IDF/Semantic/Tone = message-fit signals.",
            "ROAS range = planning scenario band.",
        ],
    )
    add_source(s, "Top Matches + Network + Text + ROI tabs")
    add_footer(s, "3 / 15")

    # 4) BOJ context
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ Context: Why This Client Case Matters")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Business context",
        [
            "BOJ is a K-beauty brand expanding in global skincare markets.",
            "Demo product: Relief Sun SPF + Glow Serum.",
            "Campaign problem: find creators with both relevance and purchase intent.",
            "Operational pain: manual creator screening is slow and inconsistent.",
        ],
        body_size=12,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Why this case is strong for demo",
        [
            "Clear product story and clear target audience.",
            "Keyword-rich skincare category for text-fit diagnostics.",
            "High sensitivity to creator-message alignment quality.",
            "Realistic benchmark comparison in U.S. skincare context.",
        ],
        body_size=12,
    )
    add_source(s, "NIQ 2025 K-beauty report + BOJ U.S. strategy references")
    add_footer(s, "4 / 15")

    # 5) BOJ run facts
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ Run Facts (Actual App Output)")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Core run stats",
        [
            f"Videos analyzed: {videos:,}",
            f"Channels scored: {channels:,}",
            f"Best ML model: {best_model}",
            f"Top-10 mean final score: {top_mean:.3f}",
            f"Top-10 mean evidence score: {ev_mean:.3f}",
            f"Top-10 community coverage: {comm_n}",
        ],
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Benchmark + ROI context",
        [
            f"Reference brand: {benchmark.get('brand', 'CeraVe')}",
            f"Reference top channel: {benchmark.get('top_channel', 'N/A')}",
            f"BOJ Top-N mean: {top_mean:.3f}",
            f"CeraVe Top-N mean: {bench_mean:.3f}",
            f"Gap (BOJ - benchmark): {bench_gap:+.3f}",
            f"ROI baseline: {float(roi.get('roas', 0.0)):.2f}x (range {float(roi.get('roas_low', 0.0)):.2f}x-{float(roi.get('roas_high', 0.0)):.2f}x)",
        ],
    )
    add_source(s, "presentation_summary_boj.json")
    add_footer(s, "5 / 15")

    # 6) top10 list + role split
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ Top-10 Recommended Influencers (Actual Output)")
    names = [f"#{i+1} {r['channel_title']} ({float(r['final_score']):.3f})" for i, (_, r) in enumerate(top.head(10).iterrows())]
    role_lines = _role_split_lines(plan)
    add_card(s, 0.68, 1.95, 5.95, 4.9, "Top-10 list", names, body_size=12)
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Role split recommendation",
        role_lines,
        body_size=11,
    )
    add_source(s, "top10_beauty_of_joseon.csv")
    add_footer(s, "6 / 15")

    # 7) representative deep dive
    s = prs.slides.add_slide(blank)
    top1_name = str(top.iloc[0]["channel_title"]) if not top.empty else "Top-1 Channel"
    add_title(s, f"Representative Deep Dive: #1 {top1_name}")
    if not top.empty:
        top1 = top.iloc[0]
        top1_plan = plan[plan["channel_title"].astype(str) == str(top1.get("channel_title", ""))]
        top1_concept = str(top1_plan.iloc[0]["content_concept"]) if not top1_plan.empty else "Concept 1 + 3 Mix"
        top1_tier = str(top1_plan.iloc[0]["activation_tier"]) if not top1_plan.empty else "Balanced"
        if "Concept 1" in top1_concept:
            backup_concept = "Concept 3 (Q&A Conversion Hook)"
        elif "Concept 2" in top1_concept:
            backup_concept = "Concept 3 (Q&A Conversion Hook)"
        else:
            backup_concept = "Concept 1 (Daily Routine)"
        latest = str(top1.get("latest_publish", "N/A"))[:10]
        profile_lines = [
            f"Final score: {float(top1.get('final_score', 0.0)):.3f}",
            f"Evidence score: {float(top1.get('evidence_score', 0.0)):.3f}",
            f"TF-IDF: {float(top1.get('tfidf_similarity', 0.0)):.3f}",
            f"SNA: {float(top1.get('sna_score', 0.0)):.3f}",
            f"n_videos used: {int(top1.get('n_videos', 0))}",
            f"Median views: {int(float(top1.get('median_views', 0) or 0)):,}",
            f"Recent publish: {latest}",
        ]
        why_lines = [
            "Why useful for BOJ:",
            "Very strong skincare-language match to campaign intent.",
            "Good context for sensitive-skin sunscreen storytelling.",
            f"Role in this run: {top1_tier}",
            "",
            "Recommended content package:",
            f"Primary: {top1_concept}",
            f"Backup: {backup_concept}",
            "Product angle: lightweight daily SPF for sensitive/acne-prone users",
        ]
    else:
        profile_lines = ["Top-1 channel data unavailable."]
        why_lines = ["No deep-dive data available."]
    add_card(s, 0.68, 1.95, 5.95, 4.9, "Channel profile (from current output)", profile_lines, body_size=12)
    add_card(s, 6.75, 1.95, 5.95, 4.9, "Business interpretation + content recommendation", why_lines, body_size=11)
    add_source(s, "top10_beauty_of_joseon.csv (row #1)")
    add_footer(s, "7 / 15")

    # 8) channel-to-content mapping
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ Channel-to-Content Strategy Mapping")
    concept_lines = _concept_mapping_lines(plan)
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Top-10 concept mapping",
        concept_lines,
        body_size=11,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "4-week execution cadence",
        [
            "Week 1: Awareness launch",
            "Week 2: Education + routine detail",
            "Week 3: Product proof/comparison",
            "Week 4: Conversion CTA + recap",
            "Then re-rank with weekly KPI feedback",
        ],
        body_size=12,
    )
    add_source(s, "Content Strategy tab + memo campaign structure")
    add_footer(s, "8 / 15")

    # 9) budget scenario
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ Budget Allocation Scenario (Current Run)")
    top_budget = plan.sort_values("plan_budget_usd", ascending=False).head(5)
    b_lines = [
        f"{r['channel_title']}: Budget ${int(r['plan_budget_usd']):,} | Conv~{int(r['plan_conversions']):,} | Revenue~${int(r['plan_revenue_usd']):,}"
        for _, r in top_budget.iterrows()
    ]
    total_budget = int(plan["plan_budget_usd"].sum())
    total_conv = int(plan["plan_conversions"].sum())
    total_rev = int(plan["plan_revenue_usd"].sum())
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Example allocation outputs",
        b_lines,
        body_size=12,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Assumption + interpretation",
        [
            "Allocation weight = 55% final score + 25% evidence + 20% reach index",
            f"Scenario budget total: ${total_budget:,}",
            f"Scenario conversions total: {total_conv:,}",
            f"Scenario revenue total: ${total_rev:,}",
            "Reach-heavy creators support awareness speed.",
            "Fit-heavy creators support conversion quality.",
        ],
        body_size=11,
    )
    add_optional_image(s, "roi_funnel_base.png", left=8.8, top=5.1, width=3.6)
    add_source(s, "presentation_summary_boj.json + derived allocation scenario")
    add_footer(s, "9 / 15")

    # 10) analysis-to-action framework
    s = prs.slides.add_slide(blank)
    add_title(s, "Analysis-to-Action Framework for BOJ")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "From analysis to decision",
        [
            "Top Matches: use as shortlist, not final spend decision alone.",
            "Evidence/Multiplier: core spend only above evidence threshold.",
            "Text Intelligence: conversion creatives only for high fit channels.",
            "Network Studio: enforce community diversity in active spend.",
            "ML Studio: validate score stability before large reallocations.",
            "ROI Lab: use scenario ranges for budget approval.",
        ],
        body_size=10,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Practical guardrails",
        [
            "Do not allocate budget by rank only.",
            "Apply pilot cap to low-activity channels.",
            "Review KPI weekly and rerank creators.",
            "Track model rationale and risk flags in exports.",
        ],
        body_size=12,
    )
    add_source(s, "Top Matches + Network + Text + ML + ROI tabs")
    add_footer(s, "10 / 15")

    # 11) 30-60-90 plan
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ 30-60-90 Day Operating Plan")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "0-30 / 31-60 days",
        [
            "0-30: Launch Top-10 with tiered budget and clean tracking links.",
            "0-30: Use 3 lead creators + 4 support creators + 3 pilots.",
            "31-60: Re-rank using CTR/CVR/ROAS.",
            "31-60: Shift budget from weak pilots to strongest converters.",
            "31-60: Update include/exclude keywords from comment feedback.",
        ],
        body_size=11,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "61-90 days",
        [
            "Build BOJ creator playbook by audience sub-segment.",
            "Segment examples: sensitive skin / acne-prone / lightweight SPF seekers.",
            "Run monthly benchmark vs competitor profile.",
            "Keep export-based audit trail for every campaign cycle.",
        ],
        body_size=12,
    )
    add_source(s, "Execution plan derived from BOJ run outputs")
    add_footer(s, "11 / 15")

    # 12) KPI governance
    s = prs.slides.add_slide(blank)
    add_title(s, "BOJ KPI Governance (Weekly)")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Core KPI stack",
        [
            "CTR",
            "CVR",
            "ROAS",
            "Branded search lift",
            "Comment quality and repeated objections",
            "Community concentration ratio",
        ],
        body_size=14,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Decision rule example",
        [
            "Increase: above-target CVR + stable evidence.",
            "Hold/test: mixed performance but strategic niche relevance.",
            "Reduce/exit: weak CVR + low evidence + off-message signals.",
            "Re-run ranking after every major budget shift.",
        ],
        body_size=12,
    )
    add_source(s, "ROI + Text + Bias diagnostics")
    add_footer(s, "12 / 15")

    # 13) industry insight pattern
    s = prs.slides.add_slide(blank)
    add_title(s, "Beauty Industry Insight #1 (from this run)")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Observed pattern",
        [
            "High reach and high semantic fit are often different creators.",
            "Single flat top-list design can waste spend.",
            "Role-split creator portfolio performs better operationally.",
        ],
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Recommended portfolio structure",
        [
            "Awareness track (reach-driven)",
            "Conversion track (fit/evidence-driven)",
            "Controlled test track (small pilot budget)",
            "Rebalance across tracks by weekly KPI",
        ],
    )
    add_source(s, "BOJ Top-10 signal pattern")
    add_footer(s, "13 / 15")

    # 14) industry operating model
    s = prs.slides.add_slide(blank)
    add_title(s, "Beauty Industry Insight #2: Scalable Operating Model")
    add_card(
        s,
        0.68,
        1.95,
        5.95,
        4.9,
        "Repeatable campaign cycle",
        [
            "Input brief -> explainable ranking -> strategy mapping -> ROI scenarios -> export",
            "Run this cycle every launch, not as one-off analysis.",
            "Use benchmark panel monthly against competitor profiles.",
        ],
        body_size=12,
    )
    add_card(
        s,
        6.75,
        1.95,
        5.95,
        4.9,
        "Governance and audit trail",
        [
            "Keep rationale text for each selected creator.",
            "Track risk flags and evidence movement over time.",
            "Store exported files for post-campaign review.",
            "Standardize framework across skincare/makeup/haircare.",
        ],
        body_size=12,
    )
    add_source(s, "AI-MCN pipeline + export architecture")
    add_footer(s, "14 / 15")

    # 15) closing
    s = prs.slides.add_slide(blank)
    add_title(s, "Closing")
    add_bullets(
        s,
        1.0,
        2.2,
        11.4,
        2.8,
        [
            "For BOJ, this output is not only a recommendation list; it is an execution plan.",
            "For beauty brands, AI-MCN turns influencer marketing into measurable decision operations.",
            "This is how we replace manual MCN-style matching with explainable AI workflow.",
        ],
        font_size=21,
    )
    thanks = s.shapes.add_textbox(Inches(3.2), Inches(5.3), Inches(7.0), Inches(0.9))
    tf = thanks.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Thank you. Q&A"
    p.font.name = "Aptos Display"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(20, 87, 148)
    p.alignment = PP_ALIGN.CENTER
    add_source(s, "AI-MCN demo output synthesis")
    add_footer(s, "15 / 15")

    prs.save(str(OUT_PPTX))
    return OUT_PPTX


if __name__ == "__main__":
    out = build()
    print(out)
